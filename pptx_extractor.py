from pptx import Presentation
from util import extractor

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    slides_data = []
    
    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip() != "":
                slide_content.append(shape.text)
        slides_data.append({
            "slide_number": slide_number,
            "content": slide_content
        })
    
    return slides_data

if __name__ == "__main__":
    input_pptx = 'test_documents/test.pptx'
    output_json = 'output_pptx.json'
    extractor(input_pptx, output_json, extract_text_from_pptx, "presentation")